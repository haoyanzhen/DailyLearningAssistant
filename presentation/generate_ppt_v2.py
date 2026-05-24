from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
C_BG      = RGBColor(0x0D, 0x0D, 0x1A)
C_DARK    = RGBColor(0x15, 0x15, 0x28)
C_SURFACE = RGBColor(0x1E, 0x1E, 0x33)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xE8, 0xE8, 0xE8)
C_GRAY    = RGBColor(0x99, 0x99, 0xAA)
C_MUTED   = RGBColor(0x66, 0x66, 0x77)
C_ACCENT  = RGBColor(0xE9, 0x6B, 0x46)   # warm orange
C_BLUE    = RGBColor(0x3A, 0x86, 0xC8)
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)
C_PURPLE  = RGBColor(0x8E, 0x44, 0xAD)
C_YELLOW  = RGBColor(0xF3, 0x9C, 0x12)
C_RED     = RGBColor(0xE7, 0x4C, 0x3C)
C_TEAL    = RGBColor(0x1A, 0xBC, 0x9C)


def bg(slide, color=C_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def circle(slide, l, t, size, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(size), Inches(size))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def label(slide, l, t, w, h, text, sz=18, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf


def para(tf, text, sz=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    p.space_before = Pt(sb)
    p.space_after = Pt(sa)
    return p


def section_header(slide, title, subtitle=""):
    label(slide, 0.8, 0.5, 10, 0.7, title, sz=34, color=C_WHITE, bold=True)
    rect(slide, 0.8, 1.25, 2.2, 0.05, C_ACCENT)
    if subtitle:
        label(slide, 0.8, 1.45, 11, 0.5, subtitle, sz=16, color=C_GRAY)


def card(slide, l, t, w, h, title, items, title_color=C_ACCENT, item_sz=13):
    rect(slide, l, t, w, h, C_SURFACE, rounded=True)
    label(slide, l + 0.25, t + 0.15, w - 0.5, 0.35, title, sz=20, color=title_color, bold=True)
    rect(slide, l + 0.25, t + 0.55, 1.2, 0.03, title_color)
    tf = label(slide, l + 0.25, t + 0.75, w - 0.5, h - 1.0, "", sz=item_sz, color=C_LIGHT)
    for item in items:
        para(tf, item, sz=item_sz, sb=5)
    return tf


def accent_line(slide, l, t, w):
    rect(slide, l, t, w, 0.05, C_ACCENT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — 封面
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)

# Decorative elements
circle(s, 10.3, 0.2, 4.5, RGBColor(0xE9, 0x6B, 0x46))
circle(s, 11.0, 3.8, 2.2, RGBColor(0x3A, 0x86, 0xC8))
circle(s, -1.5, 5.0, 3.0, RGBColor(0x8E, 0x44, 0xAD))
rect(s, 0, 0, 13.333, 0.06, C_ACCENT)

label(s, 1.2, 1.6, 10, 1.4, "每日拾光学习簿", sz=60, color=C_WHITE, bold=True)
label(s, 1.2, 3.0, 10, 0.8, "Daily Learning Assistant  v2.0", sz=32, color=C_ACCENT)
tf = label(s, 1.2, 4.4, 10, 1.0, "从 Git 提交到知识日报的全本地化 Agent 流水线", sz=22, color=C_GRAY)
para(tf, "5 个独立 Agent · Orchestrator 调度 · 拆分定时任务 · 多层回退 · LLM 追踪", sz=16, color=C_MUTED, sb=10)
label(s, 1.2, 6.3, 6, 0.4, "haoyanzhen  |  2026-05", sz=14, color=C_MUTED)

# Version badge
rect(s, 10.8, 6.2, 2.0, 0.55, C_ACCENT, rounded=True)
label(s, 10.8, 6.25, 2.0, 0.45, "v2.0", sz=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — 项目概述 & v2 新特性
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "项目概述", "一套本地运行、GitHub Pages 展示的全自动学习资产流水线")

# Left — what it is
card(s, 0.8, 2.1, 5.8, 4.8,
     "它是什么？", [
         "每日自动扫描 ~/projects 下 7 个 Git 仓库",
         "提炼数学、物理、计算机、软件工程、LLM 概念",
         "生成生动完整的中文知识讲解",
         "发布可直接阅读的 HTML 学习日报",
         "通过 SMTP 发送日报邮件通知",
         "最终形成持续积累、可检索、可回溯的个人知识库",
     ], C_ACCENT, 15)

# Right — what's new
tf = card(s, 7.2, 2.1, 5.5, 4.8,
     "v2.0 核心升级", [
         "本地独立 Agent — 不再依赖 Codex 平台",
         "Orchestrator 调度引擎 — 统一生命周期管理",
         "新增邮件 Agent — 日报自动推送到邮箱",
         "拆分调度 — HTML 生成与邮件发送独立定时",
         "多层回退 — 无变更复用 / 失败提醒 / 幂等保护",
         "状态追踪 — run_status + LLM trace 全链路可观测",
         "安全发布 — 远端检查 + 选择性暂存 + 防覆盖",
     ], C_BLUE, 14)

# Bottom tagline
rect(s, 0.8, 7.05, 11.8, 0.04, C_MUTED)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — 架构演进：Before → After
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "架构演进", "从 Codex 平台依赖 → 完全本地独立的 Agent 流水线")

# Before box
rect(s, 0.8, 2.0, 5.5, 1.0, C_SURFACE, rounded=True)
label(s, 1.0, 2.05, 5.3, 0.3, "v1.x — Codex 时代", sz=18, color=C_MUTED, bold=True)
tf = label(s, 1.0, 2.35, 5.3, 0.6, "", sz=13, color=C_GRAY)
para(tf, "依赖 Codex 平台触发 · 4 个 Agent · 无 Orchestrator · 无拆分调度 · 无回退机制", sz=13)

# Arrow
label(s, 5.8, 2.2, 1.5, 0.5, "⟶", sz=40, color=C_ACCENT, align=PP_ALIGN.CENTER)

# After box
rect(s, 7.0, 1.7, 5.5, 1.6, C_SURFACE, rounded=True)
label(s, 7.2, 1.75, 5.1, 0.3, "v2.0 — 本地独立 Agent", sz=18, color=C_ACCENT, bold=True)
tf = label(s, 7.2, 2.1, 5.1, 1.1, "", sz=13, color=C_LIGHT)
para(tf, "5 个独立 Agent + Orchestrator 调度引擎", sz=14, bold=True, sb=4)
para(tf, "launchd 本地定时 · CLI 完全可控 · 拆分调度", sz=13, sb=3)
para(tf, "多层回退 · 状态追踪 · LLM 调用链 · 安全发布", sz=13, sb=3)

# Key new modules
modules = [
    ("orchestrator/", C_ACCENT,  "调度引擎、状态管理、\n生命周期、LLM 重试"),
    ("agents/ ×5",  C_BLUE,    "5 个独立 Agent\n各司其职、CLI 契约"),
    ("scripts/",    C_GREEN,   "日报生成、邮件发送、\n配置检查、launchd 安装"),
    ("design/",     C_PURPLE,  "架构文档、回退设计、\n状态设计、任务追踪"),
]
for i, (name, color, desc) in enumerate(modules):
    x = 0.8 + i * 3.15
    rect(s, x, 3.6, 2.85, 1.6, C_SURFACE, rounded=True)
    label(s, x + 0.15, 3.7, 2.55, 0.35, name, sz=17, color=color, bold=True, align=PP_ALIGN.CENTER)
    label(s, x + 0.15, 4.1, 2.55, 0.9, desc, sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# 6 monitored repos
label(s, 0.8, 5.55, 11, 0.3, "监控范围：7 个 Git 仓库", sz=16, color=C_WHITE, bold=True)
repos = "AInote  ·  DailyLearningAssistant  ·  interview_prepare  ·  mcp  ·  ResearchPaperBase_cc  ·  ResearchPaperBase_codex"
label(s, 0.8, 5.95, 11.5, 0.4, repos, sz=13, color=C_GRAY)

rect(s, 0.8, 6.55, 11.5, 0.55, C_SURFACE, rounded=True)
label(s, 1.0, 6.62, 11, 0.4, "核心理念：Agent 负责生成 · Orchestrator 负责调度 · 中间产物可检查可复盘 · 每一步可独立补跑", sz=14, color=C_LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — 五大 Agent 流水线总览
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "五大 Agent 流水线", "每个 Agent 独立运行、CLI 契约调度、中间产物 Markdown 传递")

agents_data = [
    ("1", "代码变更\n观察员", "work_summary\n_*.md", "扫描 7 个仓库 Git 提交\n生成每日变更总结", C_ACCENT),
    ("2", "科普\n老师", "concept_relevance\n.md", "提炼 5 领域概念\n梳理关联图谱", C_BLUE),
    ("3", "知识讲解\n老师", "knowledge_explaination\n.md", "选择 3 个概念深入讲解\n维护月度知识日志", C_GREEN),
    ("4", "日报编辑\n与发布员", "learning-report\n.html", "生成 HTML 日报\n更新 manifest + Git Push", C_PURPLE),
    ("5", "邮件\n发送员", "email_preview\n.html / .txt", "生成邮件内容\nSMTP 发送日报通知", C_TEAL),
]

# Top row — agent cards
for i, (num, role, output, detail, color) in enumerate(agents_data):
    x = 0.5 + i * 2.55
    rect(s, x, 2.0, 2.3, 2.6, C_SURFACE, rounded=True)
    circle(s, x + 0.85, 2.1, 0.6, color)
    label(s, x + 0.85, 2.1, 0.6, 0.6, num, sz=22, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    label(s, x + 0.15, 2.85, 2.0, 0.55, role, sz=15, color=color, bold=True, align=PP_ALIGN.CENTER)
    label(s, x + 0.15, 3.4, 2.0, 0.45, output, sz=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < 4:
        label(s, x + 2.25, 2.85, 0.35, 0.45, "→", sz=28, color=C_GRAY, align=PP_ALIGN.CENTER)

# Bottom row — detail cards
for i, (num, role, output, detail, color) in enumerate(agents_data):
    x = 0.5 + i * 2.55
    rect(s, x, 4.95, 2.3, 2.15, C_SURFACE, rounded=True)
    label(s, x + 0.15, 5.1, 2.0, 0.3, f"Agent {num}", sz=14, color=color, bold=True)
    tf = label(s, x + 0.15, 5.5, 2.0, 1.4, "", sz=12, color=C_GRAY)
    for line in detail.split("\n"):
        para(tf, f"▸ {line}", sz=12, sb=3)

# Bottom note
label(s, 0.8, 7.1, 12, 0.3, "Orchestrator 通过 CLI 参数调度每个 Agent，不直接调用内部函数 → Agent 可独立演进，边界清晰", sz=13, color=C_MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — Agent 详解 (一)：步骤 1-3
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "Agent 详解 (一) — 内容生成阶段")

# Agent 1
card(s, 0.5, 1.9, 3.9, 5.1,
     "① 代码变更观察员", [
         "读取 ~/projects 下 7 个仓库 git log",
         "分析每次提交的文件变更与 diff 统计",
         "整理提交概览、关键文件、工作主题",
         "标注可能涉及的知识点线索",
         "无提交也生成文件并记录原因",
         "产物: work_summary_[repo].md ×7",
     ], C_ACCENT, 14)

# Agent 2
card(s, 4.7, 1.9, 3.9, 5.1,
     "② 科普老师", [
         "读取当日所有 work_summary_*.md",
         "提炼 5 领域概念：",
         "  数学 / 物理 / 计算机科学",
         "  软件工程 / 大语言模型",
         "标注来源仓库 + 概念关联图谱",
         "给出后续知识讲解建议",
         "产物: concept_relevance.md",
     ], C_BLUE, 14)

# Agent 3
card(s, 8.9, 1.9, 3.9, 5.1,
     "③ 知识讲解老师", [
         "读取 concept_relevance.md",
         "按学习价值独立选择 3 个概念",
         "不被工作相关性限制选择",
         "每个概念包含：",
         "  概念名称 · 领域 · 难度 (★)",
         "  生动完整的原理讲解",
         "  具体例子 · 常见误区",
         "  概念关联 · 开放性问题",
         "产物: knowledge_explaination.md",
         "     + knowledge_log 月度记录",
     ], C_GREEN, 13)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — Agent 详解 (二)：步骤 4-5
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "Agent 详解 (二) — 发布与分发阶段")

# Agent 4 — left
rect(s, 0.5, 1.9, 5.8, 5.1, C_SURFACE, rounded=True)
label(s, 0.8, 2.05, 5.3, 0.35, "④ 日报编辑与发布员", sz=22, color=C_PURPLE, bold=True)
rect(s, 0.8, 2.5, 1.5, 0.03, C_PURPLE)

tf = label(s, 0.8, 2.75, 5.3, 1.5, "", sz=14, color=C_LIGHT)
items_4 = [
    "读取 knowledge_explaination.md 和知识日志",
    "生成独立 HTML 学习日报页面",
    "简读 / 精读双模式切换",
    "三卡片 + 全宽文章区布局",
    "响应式：桌面并排 / 移动单列",
    "复用公共样式 style/main.css",
]
for item in items_4:
    para(tf, f"▸ {item}", sz=14, sb=5)

rect(s, 0.8, 4.8, 5.3, 1.9, RGBColor(0x25, 0x25, 0x3A), rounded=True)
label(s, 1.0, 4.9, 5.0, 0.3, "发布流程 (Orchestrator 接管)", sz=14, color=C_YELLOW, bold=True)
tf = label(s, 1.0, 5.3, 5.0, 1.2, "", sz=13, color=C_GRAY)
for item in [
    "确认 Git 工作区 + 分支匹配",
    "拉取远端 → 本地落后则停止",
    "选择性暂存 publish.paths",
    "有变更提交 / 无变更跳过",
    "本地领先则 push 到远端",
]:
    para(tf, f"→ {item}", sz=12, sb=4)

# Agent 5 — right
rect(s, 6.8, 1.9, 5.8, 5.1, C_SURFACE, rounded=True)
label(s, 7.1, 2.05, 5.3, 0.35, "⑤ 邮件发送员", sz=22, color=C_TEAL, bold=True)
rect(s, 7.1, 2.5, 1.5, 0.03, C_TEAL)

tf = label(s, 7.1, 2.75, 5.3, 3.0, "", sz=14, color=C_LIGHT)
items_5 = [
    "读取日报 manifest 和 HTML 文件",
    "LLM 生成 HTML + 纯文本邮件文案",
    "LLM 失败自动降级为静态文案",
    "默认生成预览，不发送",
    "显式 --send-email 才通过 SMTP 发送",
    "支持多收件人 + 静默列表",
    "拆分调度模式：独立定时任务",
]
for item in items_5:
    para(tf, f"▸ {item}", sz=14, sb=5)

rect(s, 7.1, 5.6, 5.3, 1.1, RGBColor(0x25, 0x25, 0x3A), rounded=True)
label(s, 7.3, 5.7, 5.0, 0.3, "失败回退", sz=14, color=C_RED, bold=True)
tf = label(s, 7.3, 5.95, 5.0, 0.6, "", sz=12, color=C_GRAY)
para(tf, "→ 日报生成失败 → 自动发送失败提醒邮件", sz=12, sb=3)
para(tf, "→ 附上上一份可用日报链接", sz=12, sb=2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Orchestrator 调度引擎
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "Orchestrator 调度引擎", "流水线的大脑 — 负责流程控制、状态管理、发布协调、回退决策")

# Core modules
mods = [
    ("run_daily.py", C_ACCENT, "总入口", "按步骤调度 Agent\n--date / --from-step / --to-step\n--only-step / --dry-run"),
    ("lifecycle.py", C_BLUE, "生命周期", "AgentSpec 能力描述\nCLI 命令构造\nLLM trace 环境注入"),
    ("state.py", C_GREEN, "状态管理", "读写 run_status.json\norchestrator_runs 拆分\n原子写入防半写"),
    ("manifest.py", C_PURPLE, "索引管理", "manifest 读写更新\n日期/路径强校验\n幂等去重 · 倒序排序"),
    ("check_status.py", C_YELLOW, "状态检查", "读取当天状态\n各 Agent 成功/失败\n诊断信息汇总"),
    ("llm.py", C_TEAL, "LLM 调用", "OpenAI Chat Completions\n失败重试 + 指数退避\nSMTP/IMAP 断联重试"),
    ("validators.py", C_RED, "结果校验", "文件存在性 · JSON 合法性\n标签协议解析\nHTML 结构校验"),
    ("config.py", C_WHITE, "配置工具", "config.json 读取\n日期/时区解析\n路径工具"),
]

for i, (name, color, title, desc) in enumerate(mods):
    col = i % 4
    row = i // 4
    x = 0.5 + col * 3.15
    y = 1.9 + row * 2.55
    rect(s, x, y, 2.9, 2.25, C_SURFACE, rounded=True)
    label(s, x + 0.2, y + 0.15, 2.5, 0.3, name, sz=14, color=color, bold=True, font="Menlo")
    label(s, x + 0.2, y + 0.45, 2.5, 0.25, title, sz=16, color=C_WHITE, bold=True)
    rect(s, x + 0.2, y + 0.75, 1.0, 0.03, color)
    tf = label(s, x + 0.2, y + 0.95, 2.5, 1.1, "", sz=12, color=C_GRAY)
    for line in desc.split("\n"):
        para(tf, f"· {line}", sz=12, sb=3)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — 拆分调度设计
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "拆分调度设计", "HTML 生成发布与邮件发送拆为两个独立 launchd 任务，时间可控、失败隔离")

# Two task boxes
tasks = [
    ("任务 A：HTML 生成 + 发布", "com.daily-learning.agent-pipeline",
     C_ACCENT, "09:00",
     ["运行 run_daily_pipeline.sh",
      "调用 run_daily.py --to-step 4 --publish",
      "执行 Agent 1-4 + Git 发布",
      "写入 orchestrator_runs.html_publish",
      "失败时停止，不进入邮件阶段"]),
    ("任务 B：邮件发送", "com.daily-learning.agent-email",
     C_TEAL, "09:30",
     ["运行 run_daily_email.sh",
      "调用 run_daily.py --only-step 5 --send-email",
      "读取 html_publish 状态判断成败",
      "成功 → 发送当天日报邮件",
      "失败 → 发送失败提醒 + 上一份日报链接"]),
]

for i, (title, service, color, time, items) in enumerate(tasks):
    x = 0.8 + i * 6.2
    rect(s, x, 1.9, 5.7, 3.5, C_SURFACE, rounded=True)
    label(s, x + 0.2, 2.0, 5.3, 0.35, title, sz=20, color=color, bold=True)
    rect(s, x + 0.2, 2.4, 1.5, 0.03, color)

    # Time badge
    rect(s, x + 4.2, 2.0, 1.3, 0.4, color, rounded=True)
    label(s, x + 4.2, 2.02, 1.3, 0.35, time, sz=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

    tf = label(s, x + 0.2, 2.65, 5.3, 2.5, "", sz=14, color=C_LIGHT)
    for item in items:
        para(tf, f"▸ {item}", sz=14, sb=6)

# Bottom — independent state keys
rect(s, 0.8, 5.7, 11.8, 1.4, C_SURFACE, rounded=True)
label(s, 1.0, 5.8, 11.4, 0.3, "状态隔离 — orchestrator_runs", sz=18, color=C_WHITE, bold=True)
tf = label(s, 1.0, 6.2, 11.4, 0.8, "", sz=13, color=C_GRAY)
para(tf, "orchestrator_runs.html_publish  —  记录 Agent 1-4 + Git 发布结果（任务 A 写入）", sz=13, sb=3)
para(tf, "orchestrator_runs.email_send      —  记录 Agent 5 邮件发送结果（任务 B 写入，并读取 A 的状态）", sz=13, sb=3)
para(tf, "两个任务写入同一 run_status.json 的不同 key，互不覆盖 → 失败隔离、独立补跑", sz=13, sb=3)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — 回退机制
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "多层回退机制", "上游缺失不编造 · 无变更智能复用 · 失败溯源可诊断 · 邮件告知不静默")

fallbacks = [
    ("无变更日回退", C_GREEN,
     "所有仓库无新增变更时",
     ["跳过 Agent 2 (概念提炼)",
      "跳过 Agent 3 (知识讲解)",
      "查找上一份有效知识讲解",
      "复制到当日目录 + 复用说明",
      "继续运行 Agent 4 生成延展型日报"]),
    ("生成阶段失败", C_RED,
     "Agent 1-4 任一步失败",
     ["停止后续生成步骤",
      "记录失败步骤 / Agent / 退出码",
      "写入诊断信息到 failures",
      "不生成基于空输入的 HTML",
      "--continue-on-failure 可覆盖"]),
    ("邮件失败提醒", C_YELLOW,
     "日报生成或发布失败",
     ["邮件任务读取 html_publish 状态",
      "失败 → 自动追加 --failure-notice",
      "发送失败提醒邮件",
      "写明失败原因",
      "附上上一份可用日报链接"]),
    ("Manifest 幂等", C_BLUE,
     "同日/同月重复运行",
     ["按日期/月份 upsert 去重",
      "不会追加重复日报条目",
      "不会追加重复知识日志",
      "保持合法 JSON + 倒序排序"]),
]

for i, (title, color, trigger, items) in enumerate(fallbacks):
    x = 0.5 + i * 3.2
    rect(s, x, 1.9, 2.95, 5.1, C_SURFACE, rounded=True)
    label(s, x + 0.15, 2.0, 2.65, 0.35, title, sz=18, color=color, bold=True)
    rect(s, x + 0.15, 2.4, 1.5, 0.03, color)
    label(s, x + 0.15, 2.6, 2.65, 0.5, f"触发：{trigger}", sz=12, color=C_GRAY)
    tf = label(s, x + 0.15, 3.2, 2.65, 3.6, "", sz=12, color=C_LIGHT)
    for item in items:
        para(tf, f"→ {item}", sz=12, sb=6)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — 状态管理与可观测性
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "状态管理与可观测性", "每一步运行可追踪 · LLM 调用可回溯 · 拆分调度状态隔离")

# Left: run_status.json structure
rect(s, 0.5, 1.9, 6.0, 5.1, C_SURFACE, rounded=True)
label(s, 0.8, 2.05, 5.5, 0.35, "run_status.json 结构", sz=20, color=C_ACCENT, bold=True)
rect(s, 0.8, 2.45, 1.5, 0.03, C_ACCENT)

tf = label(s, 0.8, 2.7, 5.5, 4.0, "", sz=12, color=C_LIGHT)
para(tf, "顶层字段：", sz=14, bold=True, sb=3)
para(tf, "  date · updated_at · agents · orchestrator", sz=12, sb=2)
para(tf, "  orchestrator_runs", sz=12, sb=2)
para(tf, "", sz=4)
para(tf, "agents.<agent_name>：", sz=14, bold=True, sb=3)
para(tf, "  status  ·  started_at / finished_at", sz=12, sb=2)
para(tf, "  input_* / output_*  文件路径", sz=12, sb=2)
para(tf, "  llm.status  ·  problems  诊断列表", sz=12, sb=2)
para(tf, "", sz=4)
para(tf, "orchestrator_runs：", sz=14, bold=True, sb=3)
para(tf, "  html_publish → Agent 1-4 + 发布", sz=12, sb=2)
para(tf, "  email_send   → Agent 5 邮件", sz=12, sb=2)
para(tf, "  各自含 status · failures · fallback_events", sz=12, sb=2)
para(tf, "  · skipped_steps · publish_events", sz=12, sb=2)

# Right: LLM trace
rect(s, 6.8, 1.9, 5.8, 5.1, C_SURFACE, rounded=True)
label(s, 7.1, 2.05, 5.3, 0.35, "LLM 追踪 (llm_trace.jsonl)", sz=20, color=C_BLUE, bold=True)
rect(s, 7.1, 2.45, 1.5, 0.03, C_BLUE)

tf = label(s, 7.1, 2.7, 5.3, 2.5, "", sz=13, color=C_LIGHT)
para(tf, "每次 LLM 调用记录一行 JSON：", sz=14, bold=True, sb=3)
para(tf, "", sz=4)
trace_fields = [
    "timestamp — 调用时间戳",
    "agent — 发起调用的 Agent 名称",
    "attempt — 第几次重试",
    "model — 使用的模型名称",
    "endpoint — API 接口地址",
    "duration_ms — 调用耗时",
    "success — 成功 / 失败",
    "error_type — 失败时的错误类型",
]
for f in trace_fields:
    para(tf, f"  · {f}", sz=12, sb=3)

para(tf, "", sz=6)
para(tf, "用途：定位 LLM 故障、评估模型性能、", sz=13, color=C_GRAY, sb=3)
para(tf, "优化 prompt 和重试策略", sz=13, color=C_GRAY, sb=2)

# Check command
rect(s, 7.1, 5.7, 5.3, 1.0, RGBColor(0x25, 0x25, 0x3A), rounded=True)
label(s, 7.3, 5.8, 5.0, 0.25, "状态检查命令", sz=14, color=C_GREEN, bold=True)
label(s, 7.3, 6.1, 5.0, 0.5, "python3 orchestrator/check_status.py --date 2026-05-23", sz=12, color=C_LIGHT, font="Menlo")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — 安全发布策略
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "安全发布策略", "无人值守场景下的 Git 发布 — 检查先行、选择性提交、防覆盖")

# Flow chart — 7 steps
steps = [
    ("1", "确认 Git\n工作区", C_MUTED),
    ("2", "检查分支\n== main", C_MUTED),
    ("3", "git fetch\n拉取远端", C_BLUE),
    ("4", "本地落后？\n→ 停止", C_RED),
    ("5", "选择性暂存\npublish.paths", C_GREEN),
    ("6", "有变更？\n→ commit", C_PURPLE),
    ("7", "领先远端？\n→ push", C_ACCENT),
]

for i, (num, label_text, color) in enumerate(steps):
    x = 0.5 + i * 1.8
    rect(s, x, 2.0, 1.55, 1.5, C_SURFACE, rounded=True)
    circle(s, x + 0.45, 2.1, 0.55, color)
    label(s, x + 0.45, 2.1, 0.55, 0.55, num, sz=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    label(s, x + 0.1, 2.8, 1.35, 0.55, label_text, sz=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
    if i < 6:
        label(s, x + 1.5, 2.45, 0.35, 0.4, "→", sz=24, color=C_GRAY, align=PP_ALIGN.CENTER)

# Key design points
points_data = [
    ("选择性暂存", C_GREEN,
     "只 git add publish.paths 中配置的生成文件，不会误提交 config.json、邮件预览等敏感/临时文件。"),
    ("远端冲突保护", C_RED,
     "发布前 git fetch 检查远端状态。如果本地落后远端（他人在此期间推送），停止发布避免覆盖。"),
    ("无需变更则跳过", C_BLUE,
     "如果暂存区无变更（所有生成文件与上次一致），跳过 git commit，避免产生空提交。"),
]

for i, (title, color, desc) in enumerate(points_data):
    y = 3.8 + i * 1.1
    rect(s, 0.5, y, 12.3, 0.95, C_SURFACE, rounded=True)
    label(s, 0.8, y + 0.1, 2.0, 0.3, title, sz=16, color=color, bold=True)
    label(s, 0.8, y + 0.5, 11.5, 0.35, desc, sz=13, color=C_GRAY)

# Publish paths example
rect(s, 0.5, 7.0, 12.3, 0.4, C_SURFACE, rounded=True)
label(s, 0.8, 7.05, 11.8, 0.3,
    'publish.paths: ["prework/...", "daily_report/...", "knowledge_log/...", "index.html"]  — 白名单机制',
    sz=11, color=C_MUTED, font="Menlo")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — 技术栈与工具链
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "技术栈与工具链")

stacks = [
    ("AI 与内容生成", C_ACCENT, [
        "LLM: DeepSeek / OpenAI 兼容接口",
        "5 个独立 Agent + Prompt Engineering",
        "Orchestrator 调度 + LLM 重试退避",
        "llm_trace.jsonl 全链路追踪",
    ]),
    ("前端展示", C_BLUE, [
        "纯静态 HTML + CSS + JavaScript",
        "单页应用模式 (index.html)",
        "CSS 公共设计系统 (style/main.css)",
        "响应式布局 (桌面 + 移动端)",
        "简读/精读渐进披露交互",
    ]),
    ("数据与状态", C_GREEN, [
        "JSON manifest 驱动导航索引",
        "Markdown 中间产物可读可查",
        "run_status.json 运行状态追踪",
        "文件系统即数据库，零依赖",
    ]),
    ("调度与部署", C_PURPLE, [
        "macOS launchd 本地定时调度",
        "拆分任务：pipeline + email",
        "Git + GitHub Pages 免费托管",
        "SSH 自动化推送",
        "零服务器成本 · 永久可访问",
    ]),
]

for i, (title, color, items) in enumerate(stacks):
    x = 0.5 + i * 3.2
    rect(s, x, 1.9, 2.95, 3.8, C_SURFACE, rounded=True)
    label(s, x + 0.2, 2.05, 2.55, 0.35, title, sz=18, color=color, bold=True)
    rect(s, x + 0.2, 2.45, 1.2, 0.03, color)
    tf = label(s, x + 0.2, 2.65, 2.55, 2.8, "", sz=13, color=C_LIGHT)
    for item in items:
        para(tf, f"▸ {item}", sz=13, sb=8)

# Bottom highlights
highlights = [
    ("Python", "核心语言"),
    ("python-pptx", "PPT 生成"),
    ("httpx", "HTTP 客户端"),
    ("launchd", "macOS 调度"),
    ("GitHub Pages", "静态托管"),
    ("SMTP/IMAP", "邮件服务"),
]
for i, (tech, desc) in enumerate(highlights):
    x = 0.5 + i * 2.1
    rect(s, x, 6.0, 1.85, 1.1, C_SURFACE, rounded=True)
    label(s, x + 0.1, 6.1, 1.65, 0.4, tech, sz=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    label(s, x + 0.1, 6.55, 1.65, 0.35, desc, sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — 总结与展望
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_header(s, "总结与展望")

# Summary
rect(s, 0.5, 1.9, 6.0, 5.0, C_SURFACE, rounded=True)
label(s, 0.8, 2.05, 5.5, 0.35, "v2.0 核心成果", sz=22, color=C_ACCENT, bold=True)
rect(s, 0.8, 2.45, 1.5, 0.03, C_ACCENT)

tf = label(s, 0.8, 2.7, 5.5, 4.0, "", sz=14, color=C_LIGHT)
summary_pts = [
    "从 Codex 依赖迁移为完全本地独立 Agent 系统",
    "5 个 Agent + Orchestrator 调度引擎",
    "拆分定时任务：生成发布 / 邮件发送 独立可控",
    "多层回退：无变更复用 · 失败停止 · 邮件通知",
    "全链路可观测：run_status + LLM trace",
    "安全发布：选择性暂存 · 远端检查 · 防覆盖",
    "CLI 完全可控：--dry-run / --from-step / --only-step",
    "单步可补跑 · 同日幂等 · manifest 去重",
    "零运维成本 · 纯静态 · GitHub Pages 托管",
]
for pt in summary_pts:
    para(tf, f"✓  {pt}", sz=14, sb=7)

# Future
rect(s, 7.0, 1.9, 5.8, 5.0, C_SURFACE, rounded=True)
label(s, 7.3, 2.05, 5.3, 0.35, "未来方向", sz=22, color=C_PURPLE, bold=True)
rect(s, 7.3, 2.45, 1.5, 0.03, C_PURPLE)

tf = label(s, 7.3, 2.7, 5.3, 4.0, "", sz=14, color=C_LIGHT)
future_pts = [
    ("跨日报全文检索", "Elasticlunr 或 Fuse.js 客户端搜索"),
    ("知识图谱可视化", "D3.js 概念关联力导向图"),
    ("RSS / Atom 订阅", "自动生成订阅源，推送每日日报"),
    ("知识回测系统", "基于历史概念自动生成练习题"),
    ("多语言讲解", "同一概念的中英双语版本"),
    ("协作知识库", "多人标注、讨论、补充概念"),
    ("Webhook 通知", "支持 Slack / 钉钉 / 飞书推送"),
    ("Docker 化部署", "跨平台一键部署"),
]
for title, desc in future_pts:
    para(tf, f"▸ {title}", sz=14, bold=True, sb=8)
    para(tf, f"     {desc}", sz=12, color=C_GRAY, sb=2)

# Bottom quote
rect(s, 0.5, 7.05, 12.3, 0.35, C_SURFACE, rounded=True)
label(s, 0.5, 7.05, 12.3, 0.35, '"把今天学到的东西，留成明天能继续走的路"', sz=16, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ── Save ──
output_path = "/Users/qingyue/projects/DailyLearningAssistant/presentation/每日拾光学习簿-v2-项目介绍.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
