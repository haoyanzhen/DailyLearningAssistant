from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
C_DARK = RGBColor(0x1A, 0x1A, 0x2E)
C_ACCENT = RGBColor(0xE9, 0x6B, 0x46)
C_ACCENT2 = RGBColor(0x3A, 0x86, 0xC8)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xF5, 0xF0, 0xEB)
C_GRAY = RGBColor(0x99, 0x99, 0x99)
C_DARKGRAY = RGBColor(0x55, 0x55, 0x55)
C_GREEN = RGBColor(0x2E, 0xCC, 0x71)
C_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
C_YELLOW = RGBColor(0xF3, 0x9C, 0x12)

def add_bg(slide, color=C_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_paragraph(tf, text, font_size=16, color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

# ==================== SLIDE 1: 封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, C_DARK)

# decorative circles
add_circle(slide, 10.5, 1.0, 3.5, RGBColor(0xE9, 0x6B, 0x46))
add_circle(slide, 11.0, 3.5, 1.8, RGBColor(0x3A, 0x86, 0xC8))

add_text_box(slide, 1.2, 1.5, 8, 1.2, "每日拾光学习簿", font_size=56, color=C_WHITE, bold=True)
add_text_box(slide, 1.2, 2.8, 10, 0.8, "Daily Learning Assistant", font_size=32, color=C_ACCENT, bold=False)

# subtitle
tf = add_text_box(slide, 1.2, 4.2, 10, 1.2, "把今天学到的东西，留成明天能继续走的路", font_size=24, color=C_GRAY)
add_paragraph(tf, "一个从 Git 提交到知识日报的全自动化学习资产流水线", font_size=18, color=C_GRAY, space_before=12)

add_text_box(slide, 1.2, 6.2, 6, 0.5, "haoyanzhen  |  2026-05", font_size=14, color=C_DARKGRAY)

# ==================== SLIDE 2: 项目概述 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "项目概述", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

# what is it
add_rounded_rect(slide, 0.8, 2.0, 5.5, 4.5, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 1.2, 2.3, 5, 0.5, "它是什么？", font_size=24, color=C_ACCENT, bold=True)
tf2 = add_text_box(slide, 1.2, 3.0, 5, 3.2, "", font_size=16, color=C_WHITE)
add_paragraph(tf2, "每日拾光学习簿是一套运行在 GitHub Pages 上的", font_size=16)
add_paragraph(tf2, "全自动学习管理系统。", font_size=16, bold=True, space_before=4)
add_paragraph(tf2, "", font_size=6)
add_paragraph(tf2, "它通过 4 个独立的 AI Agent 任务，每日自动：", font_size=16)
add_paragraph(tf2, "", font_size=4)
add_paragraph(tf2, "●  扫描 5 个 Git 仓库的提交记录", font_size=16)
add_paragraph(tf2, "●  提炼数学/物理/计算机/AI 相关概念", font_size=16)
add_paragraph(tf2, "●  生成生动详细的知识讲解", font_size=16)
add_paragraph(tf2, "●  发布可直接阅读的学习日报网页", font_size=16)
add_paragraph(tf2, "", font_size=4)
add_paragraph(tf2, "最终形成一个持续积累、可检索、可回溯的", font_size=16)
add_paragraph(tf2, "个人知识资产库。", font_size=16, bold=True)

# core features
add_rounded_rect(slide, 7.0, 2.0, 5.5, 4.5, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 7.4, 2.3, 5, 0.5, "核心特点", font_size=24, color=C_ACCENT2, bold=True)
tf2 = add_text_box(slide, 7.4, 3.0, 5, 3.2, "", font_size=16, color=C_WHITE)
items = [
    ("🤖 全自动化", "无需人工干预，每日定时触发"),
    ("📚 多领域覆盖", "数学、物理、计算机、软件工程、LLM"),
    ("🔄 流水线架构", "4 阶段接力，中间产物可检查可复盘"),
    ("🎨 渐进披露", "简读 / 精读两种阅读模式"),
    ("📊 教学记录", "月度知识索引，按日期/难度追踪"),
    ("🌐 静态站点", "GitHub Pages 托管，零服务器成本"),
]
for title, desc in items:
    p = add_paragraph(tf2, "", font_size=4, space_before=4)
    p = add_paragraph(tf2, f"{title}    {desc}", font_size=15, space_before=4)

# ==================== SLIDE 3: 监控范围 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "每日监控范围", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

add_text_box(slide, 0.8, 1.7, 11, 0.6, "每日自动扫描 ~/projects 下的 5 个 Git 仓库，生成对应工作总结", font_size=18, color=C_GRAY)

repos = [
    ("AInote", "AI 辅助开发笔记", C_ACCENT),
    ("DailyLearningAssistant", "学习助手自身", C_ACCENT2),
    ("interview_prepare", "面试准备材料", C_GREEN),
    ("ResearchPaperBase_cc", "论文管理 (Codex)", C_PURPLE),
    ("ResearchPaperBase_codex", "论文管理 (Claude Code)", C_YELLOW),
]

for i, (name, desc, color) in enumerate(repos):
    x = 0.8 + (i * 2.45)
    add_rounded_rect(slide, x, 2.7, 2.2, 3.5, RGBColor(0x25, 0x25, 0x3A))
    add_circle(slide, x + 0.75, 3.0, 0.7, color)
    add_text_box(slide, x + 0.25, 3.8, 1.7, 0.3, str(i+1), font_size=12, color=C_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, 4.1, 2.0, 0.5, name, font_size=18, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, 4.7, 2.0, 0.6, desc, font_size=12, color=C_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 6.6, 11, 0.5, "每个仓库的前一天提交 → 一份 Markdown 总结报告 → 进入概念提炼阶段", font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# ==================== SLIDE 4: 四阶段流水线概览 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "四阶段自动化流水线", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

stages = [
    ("1", "代码变更观察员", "每日工作\n总结报告", "扫描 5 个仓库的 Git 提交\n生成 work_summary_*.md", C_ACCENT),
    ("2", "科普老师", "概念提炼\n与关联", "提炼数学/物理/CS/AI 概念\n生成 concept_relevance.md", C_ACCENT2),
    ("3", "知识讲解老师", "详细知识\n讲解", "选择 3 个概念深入讲解\n生成 knowledge_explaination.md", C_GREEN),
    ("4", "日报编辑与发布员", "学习日报\n与发布", "生成 HTML 学习日报\n更新 manifest 并推送", C_PURPLE),
]

for i, (num, role, output, detail, color) in enumerate(stages):
    x = 0.5 + (i * 3.15)
    add_rounded_rect(slide, x, 1.8, 2.8, 2.3, RGBColor(0x25, 0x25, 0x3A))
    add_circle(slide, x + 1.0, 1.95, 0.5, color)
    add_text_box(slide, x + 1.0, 1.95, 0.5, 0.5, num, font_size=20, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, 2.6, 2.5, 0.4, role, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, 3.0, 2.5, 0.4, output, font_size=14, color=C_WHITE, alignment=PP_ALIGN.CENTER)

# arrows between stages
for i in range(3):
    x_arrow = 3.35 + (i * 3.15)
    add_text_box(slide, x_arrow, 2.45, 0.5, 0.5, "→", font_size=36, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# details below
for i, (num, role, output, detail, color) in enumerate(stages):
    x = 0.5 + (i * 3.15)
    add_rounded_rect(slide, x, 4.5, 2.8, 2.5, RGBColor(0x25, 0x25, 0x3A))
    add_text_box(slide, x + 0.15, 4.7, 2.5, 0.3, "产物", font_size=14, color=color, bold=True)
    add_text_box(slide, x + 0.15, 5.1, 2.5, 1.8, detail, font_size=13, color=C_GRAY)

# ==================== SLIDE 5: 阶段 1 & 2 详解 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "阶段详解 (一)", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

# Stage 1
add_rounded_rect(slide, 0.8, 1.8, 5.7, 5.0, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 1.2, 2.0, 5, 0.4, "阶段 1 — 代码变更观察员", font_size=22, color=C_ACCENT, bold=True)
tf2 = add_text_box(slide, 1.2, 2.6, 5, 4.0, "", font_size=15, color=C_WHITE)
items_s1 = [
    "●  读取前一天 ~/projects 下 5 个仓库的 git log",
    "●  分析每次提交的文件变更和 diff 统计",
    "●  整理提交概览、关键文件变更、工作主题",
    "●  标注可能涉及的知识点线索",
    "●  即使无提交也生成对应文件并记录原因",
    "●  产物: prework/YYYY-MM/YYYY-MM-DD/work_summary_*.md",
]
for item in items_s1:
    add_paragraph(tf2, item, font_size=15, space_before=6)

# Stage 2
add_rounded_rect(slide, 7.0, 1.8, 5.7, 5.0, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 7.4, 2.0, 5, 0.4, "阶段 2 — 科普老师", font_size=22, color=C_ACCENT2, bold=True)
tf2 = add_text_box(slide, 7.4, 2.6, 5, 4.0, "", font_size=15, color=C_WHITE)
items_s2 = [
    "●  读取所有当日 work_summary_*.md",
    "●  提炼 5 个领域的概念：",
    "     数学 / 物理 / 计算机科学",
    "     软件工程 / 大语言模型",
    "●  标注概念来源仓库",
    "●  梳理概念关联图谱（依赖、类比、延伸）",
    "●  给出后续知识讲解建议",
    "●  产物: concept_relevance.md",
]
for item in items_s2:
    add_paragraph(tf2, item, font_size=15, space_before=6)

# ==================== SLIDE 6: 阶段 3 & 4 详解 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "阶段详解 (二)", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

# Stage 3
add_rounded_rect(slide, 0.8, 1.8, 5.7, 5.0, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 1.2, 2.0, 5, 0.4, "阶段 3 — 知识讲解老师", font_size=22, color=C_GREEN, bold=True)
tf2 = add_text_box(slide, 1.2, 2.6, 5, 4.0, "", font_size=15, color=C_WHITE)
items_s3 = [
    "●  读取 concept_relevance.md",
    "●  按学习价值独立选择 3 个概念",
    "●  不被工作相关性限制选择",
    "●  每个概念包含:",
    "     — 概念名称与领域",
    "     — 难度评级 (★☆☆☆☆ ~ ★★★★★)",
    "     — 生动完整的原理讲解",
    "     — 具体例子与常见误区",
    "     — 概念关联与开放性问题",
    "●  产物: knowledge_explaination.md",
    "      + knowledge_log/YYYY-MM-knowledge-log.md",
]
for item in items_s3:
    add_paragraph(tf2, item, font_size=15, space_before=6)

# Stage 4
add_rounded_rect(slide, 7.0, 1.8, 5.7, 5.0, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 7.4, 2.0, 5, 0.4, "阶段 4 — 日报编辑与发布员", font_size=22, color=C_PURPLE, bold=True)
tf2 = add_text_box(slide, 7.4, 2.6, 5, 4.0, "", font_size=15, color=C_WHITE)
items_s4 = [
    "●  读取 knowledge_explaination.md",
    "●  生成独立 HTML 学习日报页面",
    "     — 简读 / 精读双模式切换",
    "     — 三卡片 + 全宽文章区布局",
    "     — 响应式：桌面并排 / 移动单列",
    "●  复用公共样式 style/main.css",
    "●  更新 manifest.json (日报+教学记录)",
    "●  检测 SSH 连通性后 git push",
    "●  产物: learning-report.html",
    "     daily_report/manifest.json",
    "     knowledge_log/manifest.json",
]
for item in items_s4:
    add_paragraph(tf2, item, font_size=15, space_before=6)

# ==================== SLIDE 7: 产物体系 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "产物体系与文件结构", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

# Left: file tree
add_rounded_rect(slide, 0.8, 1.8, 5.7, 5.0, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 1.2, 2.0, 5, 0.4, "文件结构", font_size=22, color=C_ACCENT2, bold=True)
tf2 = add_text_box(slide, 1.2, 2.6, 5, 4.0, "", font_size=14, color=C_WHITE)
tree_lines = [
    "├── prework/YYYY-MM/YYYY-MM-DD/",
    "│   ├── work_summary_AInote.md",
    "│   ├── work_summary_DailyLearningAssistant.md",
    "│   ├── work_summary_interview_prepare.md",
    "│   ├── work_summary_ResearchPaperBase_cc.md",
    "│   ├── work_summary_ResearchPaperBase_codex.md",
    "│   ├── concept_relevance.md",
    "│   └── knowledge_explaination.md",
    "├── daily_report/YYYY-MM/",
    "│   ├── YYYY-MM-DD-learning-report.html  ← 独立页面",
    "│   └── manifest.json                    ← 日报索引",
    "├── knowledge_log/",
    "│   ├── YYYY-MM-knowledge-log.md          ← 月度教学记录",
    "│   └── manifest.json                    ← 教学记录索引",
    "├── prompt/                               ← 4 个 Agent prompt",
    "├── style/                                ← 公共 CSS + JS",
    "└── index.html                            ← 首页入口",
]
for line in tree_lines:
    add_paragraph(tf2, line, font_size=14, space_before=4)

# Right: data flow
add_rounded_rect(slide, 7.0, 1.8, 5.5, 5.0, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 7.4, 2.0, 5, 0.4, "数据流向", font_size=22, color=C_PURPLE, bold=True)
tf2 = add_text_box(slide, 7.4, 2.6, 5, 4.0, "", font_size=14, color=C_WHITE)

flow_items = [
    ("Git Repos", C_ACCENT, "5 个仓库的提交记录"),
    ("work_summary_*.md", C_ACCENT2, "每日仓库变更总结"),
    ("concept_relevance.md", C_GREEN, "概念清单 + 关联图谱"),
    ("knowledge_explaination.md", C_PURPLE, "3 个概念的完整讲解"),
    ("learning-report.html", C_ACCENT, "可阅读的学习日报网页"),
    ("manifest.json", C_YELLOW, "导航索引 + 默认展示"),
]

for label, color, desc in flow_items:
    p = add_paragraph(tf2, "", font_size=4, space_before=6)
    p = add_paragraph(tf2, f"■ {label}", font_size=15, color=color, bold=True, space_before=6)
    p = add_paragraph(tf2, f"    {desc}", font_size=13, color=C_GRAY, space_before=2)
    # Add arrow (except last)
    if label != "manifest.json":
        add_paragraph(tf2, "    ↓", font_size=14, color=C_DARKGRAY, space_before=2)

# ==================== SLIDE 8: 关键设计理念 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "关键设计理念", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

principles = [
    ("中间产物契约", C_ACCENT,
     "每个阶段通过 Markdown 文件传递稳定数据，下游不猜测上游意图",
     ["固定路径规则 (YYYY-MM/YYYY-MM-DD/)", "文件名即接口约定", "出错时可独立检查任何中间文件"]),
    ("Prompt 职责边界", C_ACCENT2,
     "4 个 Agent 各有独立角色、明确输入输出和日期口径",
     ["每个 prompt 定义: 角色 + 输入 + 输出 + 约束", "降低多 Agent 上下文污染", "减少任务越界和重复工作"]),
    ("渐进披露", C_GREEN,
     "同一篇日报服务不同阅读深度 — 3 分钟速览 vs 完整深读",
     ["简读: 一句话解释 + 核心要点 + 记忆句", "精读: 完整原理 + 例子 + 误区 + 关联", "卡片式入口 → 全宽文章区展开"]),
    ("公共设计系统", C_PURPLE,
     "所有页面复用同一套 CSS 视觉语言，避免风格碎片化",
     ["style/main.css 作为唯一样式源", "日报不创建私有 CSS", "主页、归档、日报视觉一致"]),
]

y = 1.7
for title, color, desc, points in principles:
    add_rounded_rect(slide, 0.8, y, 5.5, 1.3, RGBColor(0x25, 0x25, 0x3A))
    add_text_box(slide, 1.1, y + 0.1, 5, 0.35, title, font_size=20, color=color, bold=True)
    add_text_box(slide, 1.1, y + 0.45, 5, 0.3, desc, font_size=14, color=C_GRAY)

    # detail points on right
    detail_box = add_text_box(slide, 6.8, y + 0.1, 5.5, 1.2, "", font_size=13, color=C_WHITE)
    for p in points:
        add_paragraph(detail_box, f"→ {p}", font_size=13, space_before=4)

    y += 1.45

# ==================== SLIDE 9: 近期演进 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "近期演进 (2026-05-12 ~ 05-13)", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

evolutions = [
    ("v1.0 初始发布", C_ACCENT,
     ["5 月 12 日发布首个学习日报",
      "建立 4 阶段 prompt 体系",
      "首页 + manifest 导航架构"]),
    ("v1.1 中间产物补全", C_ACCENT2,
     ["补充 work_summary 完整材料",
      "concept_relevance 概念提炼",
      "knowledge_explaination 详细讲解"]),
    ("v1.2 公共样式系统", C_GREEN,
     ["移除日报私有 CSS",
      "统一到 style/main.css",
      "主页与日报视觉一致"]),
    ("v1.3 教学记录系统", C_PURPLE,
     ["新增 knowledge_log 月度记录",
      "HTML 表格: 日期/概念/领域/难度",
      "首页 ?knowledge=1 入口"]),
    ("v1.4 阅读模式升级", C_YELLOW,
     ["简读 / 精读双模式切换",
      "三卡片 + 全宽文章区布局",
      "响应式: 桌面并排 / 移动单列"]),
]

y = 1.7
for i, (title, color, points) in enumerate(evolutions):
    x = 0.8 + (i * 2.45)
    add_rounded_rect(slide, x, y, 2.2, 5.0, RGBColor(0x25, 0x25, 0x3A))
    add_text_box(slide, x + 0.1, y + 0.15, 2, 0.4, title, font_size=17, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, x + 0.5, y + 0.6, 1.2, 0.04, color)
    tf = add_text_box(slide, x + 0.1, y + 0.8, 2, 4.0, "", font_size=12, color=C_GRAY)
    for p in points:
        add_paragraph(tf, "● " + p, font_size=12, space_before=8)

# ==================== SLIDE 10: 教学记录示例 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "教学记录概览", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

add_text_box(slide, 0.8, 1.7, 11, 0.5, "每个教学日选择 3 个概念深入讲解，累计形成月度知识索引", font_size=16, color=C_GRAY)

# Table header
col_x = [0.8, 2.8, 4.8, 7.0, 9.5]
col_w = [2.0, 2.0, 2.2, 2.5, 3.2]
headers = ["日期", "概念", "领域", "难度", "一句话解释"]
for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rounded_rect(slide, cx, 2.3, cw + 0.05, 0.5, C_ACCENT)
    add_text_box(slide, cx + 0.1, 2.35, cw, 0.4, hdr, font_size=14, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Table data rows - 2026-05-13
data_rows = [
    ("2026-05-13", "契约式开发", "软件工程 / AI 辅助开发", "★★★★", "先明确产品、接口、数据、权限和验收边界，再让实现服从约定"),
    ("", "渐进披露", "人机交互 / 信息架构", "★★★", "按用户注意力和任务阶段逐步呈现信息"),
    ("", "状态机与并发锁", "CS / 分布式系统基础", "★★★★★", "用明确状态迁移和资源占用保护，避免长任务运行态混乱"),
    ("2026-05-12", "多阶段自动化流水线", "软件工程 / Agent 工作流", "★★★★", "把复杂任务拆成职责清晰、输入输出稳定的多个自动化阶段"),
    ("", "材料一致性校验", "测试工程 / 数据质量管理", "★★★", "检查学习材料在结构、题号、语义上是否真正对应"),
    ("", "注意力机制", "深度学习 / 大语言模型", "★★★★", "让模型动态判断哪些信息更值得关注"),
]

y = 2.95
for row_idx, (date, concept, field, difficulty, explanation) in enumerate(data_rows):
    bg_color = RGBColor(0x25, 0x25, 0x3A) if row_idx % 2 == 0 else RGBColor(0x1E, 0x1E, 0x30)
    for i, (val, cx, cw) in enumerate(zip([date, concept, field, difficulty, explanation], col_x, col_w)):
        add_shape_bg(slide, cx, y, cw + 0.05, 0.55, bg_color)
        font_s = 12 if i == 4 else 13
        add_text_box(slide, cx + 0.08, y + 0.05, cw, 0.45, val, font_size=font_s, color=C_WHITE if i != 3 else C_YELLOW, alignment=PP_ALIGN.CENTER if i == 3 else PP_ALIGN.LEFT)
    y += 0.65

# ==================== SLIDE 11: 首页架构 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "首页与导航架构", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

# Three panels
panels = [
    ("最新日报", "?  (默认)", C_ACCENT,
     ["自动跳转到最新 HTML 日报",
      "显示今日知识点概览",
      "简读 / 精读切换",
      "桌面端三栏并排布局"]),
    ("学习归档", "?archive=1", C_ACCENT2,
     ["从 manifest.json 读取日报列表",
      "按年 → 月 → 日层级展开",
      "目录树浏览历史日报",
      "一键回到最新日报"]),
    ("教学记录", "?knowledge=1", C_GREEN,
     ["从 knowledge_log/manifest.json 读取",
      "按年份折叠展示月度记录",
      "每月知识点表格一览",
      "日期链接指向当天日报"]),
]

for i, (title, route, color, items) in enumerate(panels):
    x = 0.8 + (i * 4.1)
    add_rounded_rect(slide, x, 1.8, 3.8, 5.0, RGBColor(0x25, 0x25, 0x3A))
    add_text_box(slide, x + 0.2, 2.0, 3.4, 0.35, title, font_size=22, color=color, bold=True)
    add_text_box(slide, x + 0.2, 2.4, 3.4, 0.3, route, font_size=14, color=C_GRAY)
    add_shape_bg(slide, x + 0.2, 2.8, 1.5, 0.04, color)
    tf = add_text_box(slide, x + 0.2, 3.0, 3.4, 3.5, "", font_size=14, color=C_WHITE)
    for item in items:
        add_paragraph(tf, "●  " + item, font_size=14, space_before=8)

# ==================== SLIDE 12: 技术栈 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "技术栈与工具链", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

tech_groups = [
    ("内容生成", C_ACCENT, [
        "AI Agent 自动化 (Claude Code)",
        "Markdown 中间产物",
        "Prompt Engineering",
    ]),
    ("前端展示", C_ACCENT2, [
        "纯静态 HTML/CSS/JS",
        "单页应用模式 (index.html)",
        "CSS 公共设计系统",
        "响应式布局 (桌面 + 移动)",
    ]),
    ("数据管理", C_GREEN, [
        "JSON manifest 驱动导航",
        "HTML 表格教学记录",
        "文件系统即数据库",
    ]),
    ("部署托管", C_PURPLE, [
        "GitHub Pages 免费托管",
        "Git 版本控制",
        "SSH 自动化推送",
        "零服务器成本",
    ]),
]

for i, (title, color, items) in enumerate(tech_groups):
    x = 0.8 + (i * 3.15)
    add_rounded_rect(slide, x, 1.8, 2.9, 4.5, RGBColor(0x25, 0x25, 0x3A))
    add_text_box(slide, x + 0.2, 2.0, 2.5, 0.4, title, font_size=20, color=color, bold=True)
    add_shape_bg(slide, x + 0.2, 2.45, 1.5, 0.04, color)
    tf = add_text_box(slide, x + 0.2, 2.7, 2.5, 3.0, "", font_size=14, color=C_WHITE)
    for item in items:
        add_paragraph(tf, "▸  " + item, font_size=14, space_before=10)

# bottom summary
add_rounded_rect(slide, 0.8, 6.5, 11.8, 0.6, RGBColor(0x25, 0x25, 0x3A))
add_text_box(slide, 1.0, 6.55, 11.4, 0.5, "核心理念：零运维成本、全自动运行、Git 即数据库、静态页面即服务", font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# ==================== SLIDE 13: 总结与展望 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)

add_text_box(slide, 0.8, 0.5, 8, 0.7, "总结与展望", font_size=36, color=C_WHITE, bold=True)
add_shape_bg(slide, 0.8, 1.25, 2, 0.06, C_ACCENT)

# Summary
add_rounded_rect(slide, 0.8, 1.8, 5.7, 5.0, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 1.2, 2.0, 5, 0.4, "项目价值", font_size=22, color=C_ACCENT, bold=True)
tf2 = add_text_box(slide, 1.2, 2.6, 5, 4.0, "", font_size=15, color=C_WHITE)
summary_items = [
    "✓  将每日零散的代码变更转化为结构化知识",
    "",
    "✓  4 个独立 Agent 各司其职，流水线式协作",
    "",
    "✓  中间产物可检查、可复盘、可复用",
    "",
    "✓  简读/精读满足不同场景的学习需求",
    "",
    "✓  月度教学记录构建长期知识索引",
    "",
    "✓  纯静态方案：零运维、零成本、永久可访问",
]
for item in summary_items:
    add_paragraph(tf2, item, font_size=15, space_before=6)

# Future
add_rounded_rect(slide, 7.0, 1.8, 5.5, 5.0, RGBColor(0x25, 0x25, 0x3A))
tf = add_text_box(slide, 7.4, 2.0, 5, 0.4, "未来方向", font_size=22, color=C_PURPLE, bold=True)
tf2 = add_text_box(slide, 7.4, 2.6, 5, 4.0, "", font_size=15, color=C_WHITE)
future_items = [
    "→  搜索功能：跨日报全文检索",
    "",
    "→  知识图谱：概念间关联可视化",
    "",
    "→  RSS/订阅：每日推送学习日报",
    "",
    "→  知识回测：基于历史概念生成",
    "      练习题和自测",
    "",
    "→  多语言讲解：概念的多语言版本",
    "",
    "→  协作功能：多人知识库与讨论",
]
for item in future_items:
    add_paragraph(tf2, item, font_size=15, space_before=6)

# Bottom quote
add_text_box(slide, 0.8, 7.0, 12, 0.4, '"把今天学到的东西，留成明天能继续走的路"', font_size=18, color=C_ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

# Save
output_path = "/Users/qingyue/projects/DailyLearningAssistant/每日拾光学习簿-项目介绍.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
